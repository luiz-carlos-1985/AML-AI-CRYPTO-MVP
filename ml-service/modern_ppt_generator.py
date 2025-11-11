"""
Gerador de PowerPoint Ultra Moderno
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_modern_ppt(language='en'):
    """Cria PowerPoint ultra moderno"""
    
    prs = Presentation()
    
    # Configurar slide master
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    if language == 'en':
        slides_data = get_slides_data_en()
        filename = "CryptoGuard_Modern_EN.pptx"
    else:
        slides_data = get_slides_data_pt()
        filename = "CryptoGuard_Modern_PT.pptx"
    
    for slide_data in slides_data:
        add_modern_slide(prs, slide_data, slide_width, slide_height)
    
    prs.save(filename)
    return filename

def add_modern_slide(prs, slide_data, width, height):
    """Adiciona slide moderno"""
    
    # Layout em branco
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Background gradient (simulado com shape)
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, width, height
    )
    bg_fill = bg_shape.fill
    bg_fill.solid()
    
    if slide_data['type'] == 'title':
        bg_fill.fore_color.rgb = RGBColor(10, 14, 39)  # Dark blue
        add_title_content(slide, slide_data, width, height)
    elif slide_data['type'] == 'problem':
        bg_fill.fore_color.rgb = RGBColor(240, 147, 251)  # Pink gradient
        add_problem_content(slide, slide_data, width, height)
    elif slide_data['type'] == 'solution':
        bg_fill.fore_color.rgb = RGBColor(79, 172, 254)  # Blue gradient
        add_solution_content(slide, slide_data, width, height)
    else:
        bg_fill.fore_color.rgb = RGBColor(102, 126, 234)  # Purple gradient
        add_content_slide(slide, slide_data, width, height)

def add_title_content(slide, data, width, height):
    """Adiciona conteúdo do slide título"""
    
    # Título principal
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2), width - Inches(2), Inches(2)
    )
    title_frame = title_box.text_frame
    title_frame.text = data['title']
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_font = title_para.font
    title_font.name = 'Arial Black'
    title_font.size = Pt(48)
    title_font.color.rgb = RGBColor(0, 212, 255)  # Cyan
    title_font.bold = True
    
    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(4), width - Inches(2), Inches(1.5)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = data['subtitle']
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_font = subtitle_para.font
    subtitle_font.name = 'Arial'
    subtitle_font.size = Pt(24)
    subtitle_font.color.rgb = RGBColor(255, 255, 255)
    
    # Funding info
    funding_box = slide.shapes.add_textbox(
        Inches(2), Inches(6), width - Inches(4), Inches(1.5)
    )
    funding_frame = funding_box.text_frame
    funding_frame.text = data.get('funding', 'Series B Funding Round - $15M')
    funding_para = funding_frame.paragraphs[0]
    funding_para.alignment = PP_ALIGN.CENTER
    funding_font = funding_para.font
    funding_font.name = 'Arial'
    funding_font.size = Pt(20)
    funding_font.color.rgb = RGBColor(0, 255, 136)  # Green

def add_problem_content(slide, data, width, height):
    """Adiciona conteúdo do slide problema"""
    
    # Título
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), width - Inches(1), Inches(1)
    )
    title_frame = title_box.text_frame
    title_frame.text = data['title']
    title_para = title_frame.paragraphs[0]
    title_font = title_para.font
    title_font.name = 'Arial Black'
    title_font.size = Pt(36)
    title_font.color.rgb = RGBColor(255, 255, 255)
    title_font.bold = True
    
    # Estatísticas
    stats_y = Inches(2)
    for i, stat in enumerate(data['stats']):
        stat_box = slide.shapes.add_textbox(
            Inches(0.5 + i * 2.2), stats_y, Inches(2), Inches(1.5)
        )
        stat_frame = stat_box.text_frame
        
        # Número
        stat_frame.text = stat['number']
        num_para = stat_frame.paragraphs[0]
        num_para.alignment = PP_ALIGN.CENTER
        num_font = num_para.font
        num_font.name = 'Arial Black'
        num_font.size = Pt(32)
        num_font.color.rgb = RGBColor(0, 255, 136)
        num_font.bold = True
        
        # Label
        label_para = stat_frame.add_paragraph()
        label_para.text = stat['label']
        label_para.alignment = PP_ALIGN.CENTER
        label_font = label_para.font
        label_font.name = 'Arial'
        label_font.size = Pt(12)
        label_font.color.rgb = RGBColor(255, 255, 255)
    
    # Pain points
    pain_y = Inches(4.5)
    for i, pain in enumerate(data['pain_points']):
        pain_box = slide.shapes.add_textbox(
            Inches(0.5), pain_y + i * Inches(0.6), width - Inches(1), Inches(0.5)
        )
        pain_frame = pain_box.text_frame
        pain_frame.text = f"• {pain}"
        pain_para = pain_frame.paragraphs[0]
        pain_font = pain_para.font
        pain_font.name = 'Arial'
        pain_font.size = Pt(16)
        pain_font.color.rgb = RGBColor(255, 255, 255)

def add_solution_content(slide, data, width, height):
    """Adiciona conteúdo do slide solução"""
    
    # Título
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), width - Inches(1), Inches(1)
    )
    title_frame = title_box.text_frame
    title_frame.text = data['title']
    title_para = title_frame.paragraphs[0]
    title_font = title_para.font
    title_font.name = 'Arial Black'
    title_font.size = Pt(36)
    title_font.color.rgb = RGBColor(255, 255, 255)
    title_font.bold = True
    
    # Features em grid
    features_start_y = Inches(2)
    col_width = (width - Inches(2)) / 3
    
    for i, feature in enumerate(data['features']):
        col = i % 3
        row = i // 3
        
        feature_box = slide.shapes.add_textbox(
            Inches(0.5) + col * col_width,
            features_start_y + row * Inches(2),
            col_width - Inches(0.2),
            Inches(1.8)
        )
        
        feature_frame = feature_box.text_frame
        
        # Título da feature
        feature_frame.text = feature['title']
        title_para = feature_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_font = title_para.font
        title_font.name = 'Arial Black'
        title_font.size = Pt(18)
        title_font.color.rgb = RGBColor(0, 212, 255)
        title_font.bold = True
        
        # Descrição
        desc_para = feature_frame.add_paragraph()
        desc_para.text = feature['description']
        desc_para.alignment = PP_ALIGN.CENTER
        desc_font = desc_para.font
        desc_font.name = 'Arial'
        desc_font.size = Pt(12)
        desc_font.color.rgb = RGBColor(255, 255, 255)

def add_content_slide(slide, data, width, height):
    """Adiciona slide de conteúdo genérico"""
    
    # Título
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), width - Inches(1), Inches(1)
    )
    title_frame = title_box.text_frame
    title_frame.text = data['title']
    title_para = title_frame.paragraphs[0]
    title_font = title_para.font
    title_font.name = 'Arial Black'
    title_font.size = Pt(36)
    title_font.color.rgb = RGBColor(255, 255, 255)
    title_font.bold = True
    
    # Conteúdo
    content_y = Inches(2)
    for i, item in enumerate(data.get('content', [])):
        item_box = slide.shapes.add_textbox(
            Inches(0.5), content_y + i * Inches(0.7), width - Inches(1), Inches(0.6)
        )
        item_frame = item_box.text_frame
        item_frame.text = f"• {item}"
        item_para = item_frame.paragraphs[0]
        item_font = item_para.font
        item_font.name = 'Arial'
        item_font.size = Pt(18)
        item_font.color.rgb = RGBColor(255, 255, 255)

def get_slides_data_en():
    """Dados dos slides em inglês"""
    return [
        {
            'type': 'title',
            'title': 'CryptoGuard AI',
            'subtitle': 'The World\'s Most Advanced Anti-Money Laundering Platform',
            'funding': 'Series B Funding Round - $15M Investment Opportunity'
        },
        {
            'type': 'problem',
            'title': 'The $2 Trillion Money Laundering Crisis',
            'stats': [
                {'number': '$2T+', 'label': 'Money Laundered\nGlobally Per Year'},
                {'number': '$26B+', 'label': 'AML Compliance\nCosts Annually'},
                {'number': '99%+', 'label': 'Crypto Transactions\nGo Undetected'},
                {'number': '156%', 'label': 'Increase in\nCrypto Crimes'}
            ],
            'pain_points': [
                'Legacy systems can\'t handle crypto complexity',
                'Regulatory requirements constantly changing',
                'Manual processes are slow and error-prone',
                'False positives cost $25M+ per institution',
                'Cross-chain transactions are invisible'
            ]
        },
        {
            'type': 'solution',
            'title': 'CryptoGuard AI: The Ultimate AML Defense',
            'features': [
                {'title': 'Real-Time Detection', 'description': 'Sub-second analysis\nwith 95%+ accuracy'},
                {'title': 'Multi-Chain Intelligence', 'description': 'Coverage across\n8+ blockchains'},
                {'title': 'Regulatory Autopilot', 'description': 'Auto-compliance with\n50+ jurisdictions'},
                {'title': 'Enterprise Security', 'description': 'ISO 27001, SOC 2,\nPCI DSS certified'},
                {'title': 'Intelligent Reporting', 'description': 'Auto-generated\nSARs and CTRs'},
                {'title': 'Zero False Positives', 'description': '90%+ reduction\nin false alerts'}
            ]
        },
        {
            'type': 'market',
            'title': 'Massive Market Opportunity',
            'content': [
                'TAM: $45.8B - Global AML software market by 2028',
                'SAM: $12.3B - Crypto-focused AML solutions',
                'SOM: $1.2B - Enterprise crypto AML (5-year target)',
                'Crypto Growth: 2,300%+ market growth in 5 years',
                'Regulatory Pressure: $10B+ in fines annually',
                'Institutional Adoption: 87% of banks planning crypto',
                'Global Expansion: 156 countries with crypto laws'
            ]
        },
        {
            'type': 'technology',
            'title': 'Proprietary Technology Stack',
            'content': [
                'CryptoGraph Neural Network - Proprietary GNN (Patent pending)',
                'Multi-Chain Intelligence Engine - 8+ blockchain coverage',
                'Regulatory Autopilot System - 50+ jurisdiction compliance',
                'Zero-Knowledge Risk Scoring - Privacy-preserving analysis',
                'AI/ML Expertise: 50+ models, 5+ years research',
                'Data Advantage: 100M+ labeled transactions',
                'Security Infrastructure: Military-grade protection',
                'Performance Edge: Sub-second vs. competitors\' minutes'
            ]
        },
        {
            'type': 'traction',
            'title': 'Explosive Growth & Validation',
            'content': [
                'ARR: $2.3M (+340% YoY)',
                'Customers: 23 (+280% YoY)',
                'Transactions: 500M+ analyzed (+1,200% YoY)',
                'Detection: $1.2B+ money laundering detected',
                'Major US Bank: $500K+ ARR (Live in production)',
                'Top-5 Crypto Exchange: $300K+ ARR (Expanding globally)',
                'Government Agency: $200K+ ARR (Multi-year contract)',
                'Pipeline: 45+ qualified leads, $3.2M total pipeline'
            ]
        },
        {
            'type': 'competition',
            'title': 'Competitive Advantage',
            'content': [
                'Speed: Sub-second vs. competitors\' minutes/hours',
                'AI Technology: Proprietary GNN vs. basic ML',
                'Multi-Chain: 8+ blockchains vs. limited support',
                'Regulatory: Auto-compliance vs. manual processes',
                'Cost: 70% reduction vs. high implementation costs',
                'Only true next-generation AML platform',
                'Clear technology and market leadership'
            ]
        },
        {
            'type': 'funding',
            'title': 'Series B Funding Round: $15M',
            'content': [
                'Sales & Marketing (40%): $6M - Scale go-to-market globally',
                'Product Development (35%): $5.25M - Advanced AI, R&D',
                'Team Expansion (20%): $3M - Hire 50+ engineers, sales',
                'Working Capital (5%): $0.75M - Operations, legal',
                '18-Month Targets: $25M ARR, 150+ customers',
                'Geographic Expansion: 5 new markets',
                'Exit Strategy: $5B+ IPO potential in 5-7 years'
            ]
        },
        {
            'type': 'cta',
            'title': 'Join the Revolution',
            'content': [
                'Build the Future of Financial Crime Prevention',
                'Series B Investment Opportunity',
                '$15M Funding Round',
                '$75M Pre-Money Valuation',
                '10x+ Return Potential',
                'First-Mover Advantage in $45B+ Market',
                'Contact: investors@cryptoguard.ai'
            ]
        }
    ]

def get_slides_data_pt():
    """Dados dos slides em português"""
    return [
        {
            'type': 'title',
            'title': 'CryptoGuard AI',
            'subtitle': 'A Plataforma Anti-Lavagem de Dinheiro Mais Avançada do Mundo',
            'funding': 'Rodada Série B - Oportunidade de Investimento $15M'
        },
        {
            'type': 'problem',
            'title': 'A Crise de Lavagem de Dinheiro de $2 Trilhões',
            'stats': [
                {'number': '$2T+', 'label': 'Dinheiro Lavado\nGlobalmente Por Ano'},
                {'number': '$26B+', 'label': 'Custos Conformidade\nAML Anualmente'},
                {'number': '99%+', 'label': 'Transações Crypto\nNão Detectadas'},
                {'number': '156%', 'label': 'Aumento em\nCrimes Crypto'}
            ],
            'pain_points': [
                'Sistemas legados não lidam com complexidade crypto',
                'Requisitos regulatórios mudando constantemente',
                'Processos manuais lentos e propensos a erros',
                'Falsos positivos custam $25M+ por instituição',
                'Transações cross-chain são invisíveis'
            ]
        },
        {
            'type': 'solution',
            'title': 'CryptoGuard AI: A Defesa AML Definitiva',
            'features': [
                {'title': 'Detecção Tempo Real', 'description': 'Análise sub-segundo\ncom 95%+ precisão'},
                {'title': 'Inteligência Multi-Chain', 'description': 'Cobertura em\n8+ blockchains'},
                {'title': 'Piloto Automático Regulatório', 'description': 'Auto-conformidade\n50+ jurisdições'},
                {'title': 'Segurança Empresarial', 'description': 'Certificado ISO 27001\nSOC 2, PCI DSS'},
                {'title': 'Relatórios Inteligentes', 'description': 'SARs e CTRs\nauto-gerados'},
                {'title': 'Zero Falsos Positivos', 'description': '90%+ redução\nalertas falsos'}
            ]
        },
        {
            'type': 'market',
            'title': 'Oportunidade Massiva de Mercado',
            'content': [
                'TAM: $45.8B - Mercado AML global até 2028',
                'SAM: $12.3B - Soluções AML focadas em crypto',
                'SOM: $1.2B - Meta AML crypto empresarial (5 anos)',
                'Crescimento Crypto: 2,300%+ crescimento em 5 anos',
                'Pressão Regulatória: $10B+ multas anualmente',
                'Adoção Institucional: 87% bancos planejando crypto',
                'Expansão Global: 156 países com leis crypto'
            ]
        },
        {
            'type': 'technology',
            'title': 'Stack Tecnológico Proprietário',
            'content': [
                'Rede Neural CryptoGraph - GNN proprietária (Patente pendente)',
                'Motor Inteligência Multi-Chain - Cobertura 8+ blockchains',
                'Sistema Piloto Automático Regulatório - Conformidade 50+ jurisdições',
                'Pontuação Risco Zero-Knowledge - Análise preservando privacidade',
                'Expertise AI/ML: 50+ modelos, 5+ anos pesquisa',
                'Vantagem Dados: 100M+ transações rotuladas',
                'Infraestrutura Segurança: Proteção grau militar',
                'Vantagem Performance: Sub-segundo vs. minutos concorrentes'
            ]
        },
        {
            'type': 'traction',
            'title': 'Crescimento Explosivo e Validação',
            'content': [
                'ARR: $2.3M (+340% YoY)',
                'Clientes: 23 (+280% YoY)',
                'Transações: 500M+ analisadas (+1,200% YoY)',
                'Detecção: $1.2B+ lavagem dinheiro detectada',
                'Grande Banco Americano: $500K+ ARR (Produção ativa)',
                'Top-5 Exchange Crypto: $300K+ ARR (Expandindo globalmente)',
                'Agência Governamental: $200K+ ARR (Contrato multi-ano)',
                'Pipeline: 45+ leads qualificados, $3.2M pipeline total'
            ]
        },
        {
            'type': 'competition',
            'title': 'Vantagem Competitiva',
            'content': [
                'Velocidade: Sub-segundo vs. minutos/horas concorrentes',
                'Tecnologia IA: GNN proprietária vs. ML básico',
                'Multi-Chain: 8+ blockchains vs. suporte limitado',
                'Regulatório: Auto-conformidade vs. processos manuais',
                'Custo: 70% redução vs. altos custos implementação',
                'Única plataforma AML verdadeiramente próxima geração',
                'Liderança tecnológica e de mercado clara'
            ]
        },
        {
            'type': 'funding',
            'title': 'Rodada Série B: $15M',
            'content': [
                'Vendas & Marketing (40%): $6M - Escalar go-to-market globalmente',
                'Desenvolvimento Produto (35%): $5.25M - IA avançada, P&D',
                'Expansão Equipe (20%): $3M - Contratar 50+ engenheiros, vendas',
                'Capital Giro (5%): $0.75M - Operações, jurídico',
                'Metas 18 Meses: $25M ARR, 150+ clientes',
                'Expansão Geográfica: 5 novos mercados',
                'Estratégia Saída: Potencial IPO $5B+ em 5-7 anos'
            ]
        },
        {
            'type': 'cta',
            'title': 'Junte-se à Revolução',
            'content': [
                'Construa o Futuro da Prevenção ao Crime Financeiro',
                'Oportunidade de Investimento Série B',
                'Rodada de Financiamento $15M',
                'Valuation Pré-Money $75M',
                'Potencial Retorno 10x+',
                'Vantagem First-Mover em Mercado $45B+',
                'Contato: investors@cryptoguard.ai'
            ]
        }
    ]

if __name__ == "__main__":
    print("Gerando PowerPoints ultra modernos...")
    
    ppt_en = create_modern_ppt('en')
    ppt_pt = create_modern_ppt('pt')
    
    print(f"\nPowerPoints modernos criados:")
    print(f"Inglês: {ppt_en}")
    print(f"Português: {ppt_pt}")
    print(f"\nCaracterísticas:")
    print("- Design moderno com cores vibrantes")
    print("- Layouts profissionais")
    print("- Tipografia impactante")
    print("- Backgrounds gradientes")
    print("- Formatação otimizada")
    print(f"\nAbra os arquivos no PowerPoint para apresentar!")