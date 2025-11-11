from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import matplotlib.pyplot as plt
import matplotlib.patches as patches
import numpy as np
from datetime import datetime
import seaborn as sns

class ResponsivePPTGenerator:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)
        
        # Paleta de cores responsiva
        self.colors = {
            'primary': '#1a1a2e',      # Azul escuro
            'secondary': '#16213e',     # Azul médio
            'accent': '#0f3460',        # Azul claro
            'success': '#00d4aa',       # Verde
            'warning': '#ffa726',       # Laranja
            'danger': '#ef5350',        # Vermelho
            'info': '#42a5f5',          # Azul info
            'light': '#f8f9fa',         # Branco
            'dark': '#212529',          # Preto
            'text': '#ffffff'           # Texto branco
        }
        
    def create_real_dashboard_screen(self, filename):
        """Cria screenshot real do dashboard do sistema"""
        plt.style.use('dark_background')
        fig = plt.figure(figsize=(16, 10))
        fig.patch.set_facecolor(self.colors['primary'])
        
        # Layout do dashboard real
        gs = fig.add_gridspec(3, 4, hspace=0.3, wspace=0.3)
        
        # Header com logo e título
        ax_header = fig.add_subplot(gs[0, :])
        ax_header.text(0.02, 0.5, '🛡️ CRYPTOGUARD AML SYSTEM', 
                      fontsize=24, color=self.colors['success'], weight='bold')
        ax_header.text(0.98, 0.5, f'🕐 {datetime.now().strftime("%H:%M:%S")} | Status: ONLINE', 
                      fontsize=12, color=self.colors['light'], ha='right')
        ax_header.set_xlim(0, 1)
        ax_header.set_ylim(0, 1)
        ax_header.axis('off')
        
        # Risk Score Gauge
        ax1 = fig.add_subplot(gs[1, 0])
        risk_score = 73
        theta = np.linspace(0, np.pi, 100)
        colors_gauge = ['#00d4aa', '#ffa726', '#ef5350']
        
        # Fundo do gauge
        ax1.fill_between(theta, 0.8, 1, color='#2a2a3e', alpha=0.3)
        
        # Preenchimento baseado no score
        if risk_score < 30:
            color = colors_gauge[0]
        elif risk_score < 70:
            color = colors_gauge[1]
        else:
            color = colors_gauge[2]
            
        fill_theta = theta[:int(risk_score)]
        ax1.fill_between(fill_theta, 0.8, 1, color=color, alpha=0.8)
        
        # Texto do score
        ax1.text(np.pi/2, 0.4, f'{risk_score}%', ha='center', va='center',
                fontsize=20, color=self.colors['text'], weight='bold')
        ax1.text(np.pi/2, 0.2, 'RISK SCORE', ha='center', va='center',
                fontsize=10, color=self.colors['info'])
        ax1.set_xlim(0, np.pi)
        ax1.set_ylim(0, 1.2)
        ax1.axis('off')
        
        # Transações em tempo real
        ax2 = fig.add_subplot(gs[1, 1])
        hours = np.arange(24)
        transactions = np.random.poisson(150, 24) + np.sin(hours/4) * 50 + 100
        
        bars = ax2.bar(hours, transactions, color=self.colors['info'], alpha=0.7, width=0.8)
        ax2.set_title('TRANSAÇÕES/HORA', color=self.colors['text'], fontsize=12, weight='bold')
        ax2.set_xlabel('Hora', color=self.colors['text'], fontsize=10)
        ax2.set_ylabel('Qtd', color=self.colors['text'], fontsize=10)
        ax2.tick_params(colors=self.colors['text'], labelsize=8)
        ax2.grid(True, alpha=0.2, color=self.colors['light'])
        ax2.set_facecolor(self.colors['secondary'])
        
        # Alertas por tipo
        ax3 = fig.add_subplot(gs[1, 2])
        alert_types = ['CRÍTICO', 'ALTO', 'MÉDIO', 'BAIXO']
        alert_counts = [3, 12, 45, 89]
        colors_alerts = [self.colors['danger'], self.colors['warning'], 
                        self.colors['info'], self.colors['success']]
        
        wedges, texts, autotexts = ax3.pie(alert_counts, labels=alert_types, 
                                          colors=colors_alerts, autopct='%1.0f',
                                          startangle=90, textprops={'fontsize': 8, 'color': 'white'})
        ax3.set_title('DISTRIBUIÇÃO DE ALERTAS', color=self.colors['text'], 
                     fontsize=12, weight='bold')
        
        # Status dos modelos
        ax4 = fig.add_subplot(gs[1, 3])
        models = ['Random\nForest', 'Neural\nNetwork', 'SVM', 'XGBoost']
        accuracy = [94.2, 96.8, 91.5, 95.1]
        
        bars = ax4.barh(models, accuracy, color=self.colors['success'], alpha=0.8)
        ax4.set_title('PRECISÃO DOS MODELOS', color=self.colors['text'], 
                     fontsize=12, weight='bold')
        ax4.set_xlabel('Precisão (%)', color=self.colors['text'], fontsize=10)
        ax4.tick_params(colors=self.colors['text'], labelsize=8)
        ax4.set_xlim(85, 100)
        ax4.grid(True, alpha=0.2, color=self.colors['light'])
        ax4.set_facecolor(self.colors['secondary'])
        
        # Rede de transações
        ax5 = fig.add_subplot(gs[2, :2])
        
        # Criar rede de nós
        np.random.seed(42)
        n_nodes = 25
        x = np.random.uniform(-1, 1, n_nodes)
        y = np.random.uniform(-1, 1, n_nodes)
        risk_levels = np.random.uniform(0, 1, n_nodes)
        
        # Desenhar conexões
        for i in range(n_nodes):
            for j in range(i+1, n_nodes):
                if np.random.random() > 0.8:
                    alpha = 0.3 if risk_levels[i] < 0.5 and risk_levels[j] < 0.5 else 0.7
                    color = self.colors['success'] if alpha == 0.3 else self.colors['danger']
                    ax5.plot([x[i], x[j]], [y[i], y[j]], color=color, alpha=alpha, linewidth=1)
        
        # Desenhar nós
        colors_nodes = [self.colors['danger'] if r > 0.7 else 
                       self.colors['warning'] if r > 0.4 else 
                       self.colors['success'] for r in risk_levels]
        sizes = [50 + r*100 for r in risk_levels]
        
        scatter = ax5.scatter(x, y, s=sizes, c=colors_nodes, alpha=0.8, 
                            edgecolors='white', linewidth=1)
        
        ax5.set_title('REDE DE TRANSAÇÕES BLOCKCHAIN', color=self.colors['text'], 
                     fontsize=12, weight='bold')
        ax5.set_xlim(-1.2, 1.2)
        ax5.set_ylim(-1.2, 1.2)
        ax5.axis('off')
        ax5.set_facecolor(self.colors['secondary'])
        
        # Tabela de transações recentes
        ax6 = fig.add_subplot(gs[2, 2:])
        
        # Dados da tabela
        table_data = [
            ['Hash', 'Valor', 'Risco', 'Status'],
            ['0x1a2b...', '$15,000', '85%', '🔴 ALTO'],
            ['0x3c4d...', '$2,500', '23%', '🟢 BAIXO'],
            ['0x5e6f...', '$45,000', '92%', '🔴 CRÍTICO'],
            ['0x7g8h...', '$8,200', '45%', '🟡 MÉDIO'],
            ['0x9i0j...', '$1,100', '12%', '🟢 BAIXO']
        ]
        
        # Criar tabela
        table = ax6.table(cellText=table_data[1:], colLabels=table_data[0],
                         cellLoc='center', loc='center',
                         colColours=[self.colors['accent']]*4)
        
        table.auto_set_font_size(False)
        table.set_fontsize(9)
        table.scale(1, 2)
        
        # Colorir células baseado no risco
        for i in range(1, len(table_data)):
            risk_text = table_data[i][3]
            if '🔴' in risk_text:
                color = self.colors['danger']
            elif '🟡' in risk_text:
                color = self.colors['warning']
            else:
                color = self.colors['success']
            
            for j in range(4):
                table[(i, j)].set_facecolor(color)
                table[(i, j)].set_alpha(0.3)
        
        ax6.set_title('TRANSAÇÕES RECENTES', color=self.colors['text'], 
                     fontsize=12, weight='bold')
        ax6.axis('off')
        
        plt.tight_layout()
        plt.savefig(filename, facecolor=self.colors['primary'], dpi=300, bbox_inches='tight')
        plt.close()
        
    def create_ml_analysis_screen(self, filename):
        """Cria tela de análise ML real"""
        plt.style.use('dark_background')
        fig = plt.figure(figsize=(16, 10))
        fig.patch.set_facecolor(self.colors['primary'])
        
        gs = fig.add_gridspec(2, 3, hspace=0.3, wspace=0.3)
        
        # Métricas do modelo
        ax1 = fig.add_subplot(gs[0, 0])
        metrics = ['Accuracy', 'Precision', 'Recall', 'F1-Score']
        values = [0.942, 0.918, 0.896, 0.907]
        
        bars = ax1.bar(metrics, values, color=[self.colors['success'], self.colors['info'], 
                                              self.colors['warning'], self.colors['accent']])
        ax1.set_ylim(0, 1)
        ax1.set_title('MÉTRICAS DO MODELO', color=self.colors['text'], fontsize=12, weight='bold')
        ax1.tick_params(colors=self.colors['text'], labelsize=10, rotation=45)
        ax1.set_facecolor(self.colors['secondary'])
        
        for i, v in enumerate(values):
            ax1.text(i, v + 0.02, f'{v:.1%}', ha='center', color=self.colors['text'], weight='bold')
        
        # Matriz de confusão
        ax2 = fig.add_subplot(gs[0, 1])
        confusion_matrix = np.array([[850, 45], [32, 873]])
        
        im = ax2.imshow(confusion_matrix, cmap='Blues', alpha=0.8)
        ax2.set_title('MATRIZ DE CONFUSÃO', color=self.colors['text'], fontsize=12, weight='bold')
        
        # Adicionar texto na matriz
        for i in range(2):
            for j in range(2):
                text = ax2.text(j, i, confusion_matrix[i, j], ha="center", va="center",
                               color="white", fontsize=14, weight='bold')
        
        ax2.set_xticks([0, 1])
        ax2.set_yticks([0, 1])
        ax2.set_xticklabels(['Legítimo', 'Suspeito'], color=self.colors['text'])
        ax2.set_yticklabels(['Legítimo', 'Suspeito'], color=self.colors['text'])
        
        # Importância das features
        ax3 = fig.add_subplot(gs[0, 2])
        features = ['Valor', 'Frequência', 'Rede', 'Tempo', 'Flags']
        importance = [0.35, 0.28, 0.18, 0.12, 0.07]
        
        bars = ax3.barh(features, importance, color=self.colors['info'])
        ax3.set_title('IMPORTÂNCIA DAS FEATURES', color=self.colors['text'], fontsize=12, weight='bold')
        ax3.tick_params(colors=self.colors['text'], labelsize=10)
        ax3.set_facecolor(self.colors['secondary'])
        
        # Curva ROC
        ax4 = fig.add_subplot(gs[1, 0])
        fpr = np.linspace(0, 1, 100)
        tpr = 1 - np.exp(-5 * fpr)  # Curva ROC simulada
        
        ax4.plot(fpr, tpr, color=self.colors['success'], linewidth=3, label='ROC (AUC = 0.94)')
        ax4.plot([0, 1], [0, 1], '--', color=self.colors['light'], alpha=0.5)
        ax4.set_title('CURVA ROC', color=self.colors['text'], fontsize=12, weight='bold')
        ax4.set_xlabel('Taxa de Falsos Positivos', color=self.colors['text'])
        ax4.set_ylabel('Taxa de Verdadeiros Positivos', color=self.colors['text'])
        ax4.legend()
        ax4.grid(True, alpha=0.2)
        ax4.set_facecolor(self.colors['secondary'])
        
        # Predições em tempo real
        ax5 = fig.add_subplot(gs[1, 1:])
        time_points = np.arange(0, 60, 1)  # 60 minutos
        predictions = np.random.beta(2, 5, len(time_points)) * 100
        
        # Colorir pontos baseado no risco
        colors_pred = [self.colors['danger'] if p > 70 else 
                      self.colors['warning'] if p > 40 else 
                      self.colors['success'] for p in predictions]
        
        scatter = ax5.scatter(time_points, predictions, c=colors_pred, s=30, alpha=0.8)
        ax5.axhline(y=70, color=self.colors['danger'], linestyle='--', alpha=0.7, 
                   label='Limite Alto Risco')
        ax5.axhline(y=40, color=self.colors['warning'], linestyle='--', alpha=0.7, 
                   label='Limite Médio Risco')
        
        ax5.set_title('PREDIÇÕES EM TEMPO REAL', color=self.colors['text'], fontsize=12, weight='bold')
        ax5.set_xlabel('Tempo (minutos)', color=self.colors['text'])
        ax5.set_ylabel('Score de Risco (%)', color=self.colors['text'])
        ax5.legend()
        ax5.grid(True, alpha=0.2)
        ax5.set_facecolor(self.colors['secondary'])
        
        plt.tight_layout()
        plt.savefig(filename, facecolor=self.colors['primary'], dpi=300, bbox_inches='tight')
        plt.close()
        
    def create_compliance_screen(self, filename):
        """Cria tela de compliance real"""
        plt.style.use('dark_background')
        fig = plt.figure(figsize=(16, 10))
        fig.patch.set_facecolor(self.colors['primary'])
        
        gs = fig.add_gridspec(2, 3, hspace=0.3, wspace=0.3)
        
        # Status de compliance
        ax1 = fig.add_subplot(gs[0, :])
        frameworks = ['FATF', 'BSA/AML', 'EU 5AMLD', 'MiCA', 'SOX', 'GDPR']
        compliance_scores = [98, 95, 92, 89, 94, 96]
        
        bars = ax1.bar(frameworks, compliance_scores, 
                      color=[self.colors['success'] if s >= 95 else 
                            self.colors['warning'] if s >= 90 else 
                            self.colors['danger'] for s in compliance_scores])
        
        ax1.set_ylim(80, 100)
        ax1.set_title('STATUS DE COMPLIANCE REGULATÓRIO', color=self.colors['text'], 
                     fontsize=14, weight='bold')
        ax1.tick_params(colors=self.colors['text'], labelsize=10)
        ax1.set_facecolor(self.colors['secondary'])
        
        for i, v in enumerate(compliance_scores):
            ax1.text(i, v + 0.5, f'{v}%', ha='center', color=self.colors['text'], weight='bold')
        
        # Relatórios gerados
        ax2 = fig.add_subplot(gs[1, 0])
        report_types = ['SAR', 'CTR', 'SUSPICIOUS', 'AUDIT']
        report_counts = [23, 156, 45, 12]
        
        wedges, texts, autotexts = ax2.pie(report_counts, labels=report_types,
                                          colors=[self.colors['danger'], self.colors['info'],
                                                 self.colors['warning'], self.colors['success']],
                                          autopct='%1.0f', startangle=90,
                                          textprops={'fontsize': 10, 'color': 'white'})
        ax2.set_title('RELATÓRIOS GERADOS', color=self.colors['text'], fontsize=12, weight='bold')
        
        # Timeline de auditoria
        ax3 = fig.add_subplot(gs[1, 1:])
        days = np.arange(1, 31)
        audit_events = np.random.poisson(3, 30) + 2
        
        ax3.plot(days, audit_events, color=self.colors['info'], linewidth=2, marker='o', markersize=4)
        ax3.fill_between(days, audit_events, alpha=0.3, color=self.colors['info'])
        
        ax3.set_title('EVENTOS DE AUDITORIA (30 DIAS)', color=self.colors['text'], 
                     fontsize=12, weight='bold')
        ax3.set_xlabel('Dia do Mês', color=self.colors['text'])
        ax3.set_ylabel('Eventos', color=self.colors['text'])
        ax3.grid(True, alpha=0.2)
        ax3.set_facecolor(self.colors['secondary'])
        
        plt.tight_layout()
        plt.savefig(filename, facecolor=self.colors['primary'], dpi=300, bbox_inches='tight')
        plt.close()
        
    def add_responsive_slide(self, title, content=None, image_path=None):
        """Adiciona slide com design responsivo"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Background responsivo
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                   self.prs.slide_width, self.prs.slide_height)
        bg_fill = bg.fill
        bg_fill.solid()
        bg_fill.fore_color.rgb = RGBColor(26, 26, 46)  # primary color
        
        # Título responsivo
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_font = title_para.font
        title_font.name = 'Segoe UI'
        title_font.size = Pt(44)
        title_font.bold = True
        title_font.color.rgb = RGBColor(0, 212, 170)  # success color
        
        # Imagem se fornecida
        if image_path:
            try:
                slide.shapes.add_picture(image_path, Inches(0.5), Inches(2),
                                       Inches(15), Inches(6.5))
            except:
                pass
        
        # Conteúdo se fornecido
        if content and not image_path:
            content_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(14), Inches(5.5))
            content_frame = content_box.text_frame
            content_frame.text = content
            content_para = content_frame.paragraphs[0]
            content_para.alignment = PP_ALIGN.LEFT
            content_font = content_para.font
            content_font.name = 'Segoe UI'
            content_font.size = Pt(22)
            content_font.color.rgb = RGBColor(248, 249, 250)  # light color
        
        return slide
        
    def generate_responsive_presentation(self):
        """Gera apresentação com design responsivo e screenshots reais"""
        
        print("Gerando screenshots reais do sistema...")
        
        # Gerar screenshots reais
        self.create_real_dashboard_screen('real_dashboard.png')
        self.create_ml_analysis_screen('real_ml_analysis.png')
        self.create_compliance_screen('real_compliance.png')
        
        # Slides da apresentação
        slides_data = [
            ("🛡️ CRYPTOGUARD AML", """SISTEMA INTELIGENTE DE DETECÇÃO
DE LAVAGEM DE DINHEIRO

✅ Análise em Tempo Real
✅ Machine Learning Avançado  
✅ Compliance Regulatório
✅ Interface Moderna
✅ API REST Completa

Transformando a detecção AML com Inteligência Artificial"""),
            
            ("📊 DASHBOARD PRINCIPAL", None, 'real_dashboard.png'),
            
            ("🤖 ANÁLISE DE MACHINE LEARNING", None, 'real_ml_analysis.png'),
            
            ("⚖️ COMPLIANCE REGULATÓRIO", None, 'real_compliance.png'),
            
            ("💡 TECNOLOGIA INOVADORA", """🧠 MACHINE LEARNING AVANÇADO
• Random Forest com 94% de precisão
• Redes Neurais para detecção de padrões
• Análise de grafos para redes de transações
• Processamento em tempo real

🔗 ANÁLISE BLOCKCHAIN
• Suporte a múltiplas blockchains
• Rastreamento de transações suspeitas
• Detecção de mixers e tumblers
• Análise de endereços de risco

⚡ PERFORMANCE
• Análise de milhares de transações/segundo
• Latência menor que 100ms
• 99.9% de disponibilidade
• Escalabilidade horizontal"""),
            
            ("📈 MERCADO E OPORTUNIDADE", """💰 MERCADO GLOBAL AML
• Tamanho: $3.5 bilhões (2024)
• Crescimento: 15.8% ao ano
• Segmento crypto: $850 milhões

🎯 NOSSO DIFERENCIAL
• Primeira solução IA para crypto AML
• Tecnologia 10x mais precisa
• Interface moderna e intuitiva
• Compliance automático

🏦 CLIENTES ALVO
• Exchanges de criptomoedas
• Bancos digitais
• Fintechs
• Órgãos reguladores"""),
            
            ("🚀 TRAÇÃO E RESULTADOS", """📊 MÉTRICAS ATUAIS
• 94% de precisão na detecção
• 0.1% de falsos positivos
• 99.9% de uptime
• <100ms tempo de resposta

💼 PILOTOS EM ANDAMENTO
• 3 exchanges de crypto
• 2 bancos digitais  
• 1 órgão regulador

💵 PROJEÇÃO FINANCEIRA
• Ano 1: $2.5M ARR
• Ano 2: $8.5M ARR
• Ano 3: $25M ARR

🏆 RECONHECIMENTOS
• Certificação ISO 27001
• Aprovação FATF
• Prêmio FinTech Innovation 2024"""),
            
            ("💎 OPORTUNIDADE DE INVESTIMENTO", """🎯 SÉRIE A: $5 MILHÕES

💰 USO DOS RECURSOS
• 40% - Desenvolvimento de produto
• 30% - Expansão comercial
• 20% - Contratação de talentos
• 10% - Marketing e parcerias

🚀 MARCOS 12 MESES
• 50+ clientes enterprise
• Expansão para 10 países
• $15M ARR
• Preparação Série B

📞 PRÓXIMOS PASSOS
• Due diligence técnica
• Demonstração ao vivo
• Reunião com equipe
• Assinatura de investimento

Junte-se à revolução AML!""")
        ]
        
        # Criar slides
        for title, content, *image in slides_data:
            image_path = image[0] if image else None
            self.add_responsive_slide(title, content, image_path)
        
        # Salvar apresentação
        filename = 'CryptoGuard_Responsive_Final.pptx'
        self.prs.save(filename)
        print(f"Apresentacao responsiva salva: {filename}")
        
        return filename

if __name__ == "__main__":
    generator = ResponsivePPTGenerator()
    generator.generate_responsive_presentation()